"""
app/extraction/pptx_extractor.py

Multimodal PPTX extraction: slide titles (as headings), body text, tables,
images, OCR text inside images, captions, and speaker notes -- linked
together as Blocks (see app/extraction/base.py). Reading order is
slide-order + on-slide shape order (python-pptx exposes shapes in the order
they were added, which is a reasonable reading-order proxy for slides).

Does not touch app/ingestion/pipeline.py.
"""

from __future__ import annotations

import os
import uuid

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract

from app.extraction.base import (
    BaseExtractor,
    Block,
    BlockType,
    ExtractionResult,
)

CAPTION_PREFIXES = ("figure", "fig.", "fig ", "table", "chart", "diagram")
CAPTION_LINK_WINDOW = 2


class PptxExtractor(BaseExtractor):
    file_types = ("pptx",)

    def __init__(self, image_output_dir: str = "extracted_images"):
        self.image_output_dir = image_output_dir

    def extract(self, file_path: str, doc_id: str) -> ExtractionResult:
        result = ExtractionResult(
            doc_id=doc_id,
            source_file=os.path.basename(file_path),
            file_type="pptx",
        )

        doc_image_dir = os.path.join(self.image_output_dir, doc_id)
        os.makedirs(doc_image_dir, exist_ok=True)

        try:
            prs = Presentation(file_path)
        except Exception as e:
            result.warnings.append(f"Failed to open PPTX: {e}")
            return result

        order_index = 0
        image_counter = 0

        for slide_number, slide in enumerate(prs.slides):
            title_shape = slide.shapes.title

            for shape in slide.shapes:
                is_title = title_shape is not None and shape.shape_id == title_shape.shape_id

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_counter += 1
                        image = shape.image
                        image_filename = f"slide{slide_number}_img{image_counter}.{image.ext}"
                        image_path = os.path.join(doc_image_dir, image_filename)
                        with open(image_path, "wb") as f:
                            f.write(image.blob)

                        ocr_text = self._run_ocr(image_path, result.warnings)

                        image_block = Block(
                            block_id=str(uuid.uuid4()),
                            doc_id=doc_id,
                            source_file=result.source_file,
                            block_type=BlockType.IMAGE,
                            content="",
                            image_path=image_path,
                            page_number=slide_number,
                            order_index=order_index,
                        )
                        result.blocks.append(image_block)
                        order_index += 1

                        if ocr_text.strip():
                            ocr_block = Block(
                                block_id=str(uuid.uuid4()),
                                doc_id=doc_id,
                                source_file=result.source_file,
                                block_type=BlockType.OCR_TEXT,
                                content=ocr_text,
                                page_number=slide_number,
                                order_index=order_index,
                                relates_to=[image_block.block_id],
                            )
                            image_block.relates_to.append(ocr_block.block_id)
                            result.blocks.append(ocr_block)
                            order_index += 1

                    except Exception as e:
                        result.warnings.append(
                            f"Image extraction failed (slide {slide_number}): {e}"
                        )
                    continue

                if shape.has_table:
                    markdown, rows = self._table_to_markdown(shape.table)
                    if markdown.strip():
                        result.blocks.append(
                            Block(
                                block_id=str(uuid.uuid4()),
                                doc_id=doc_id,
                                source_file=result.source_file,
                                block_type=BlockType.TABLE,
                                content=markdown,
                                structured_data=rows,
                                page_number=slide_number,
                                order_index=order_index,
                            )
                        )
                        order_index += 1
                    continue

                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue

                    is_caption = text.lower().startswith(CAPTION_PREFIXES)
                    block_type = (
                        BlockType.HEADING if is_title
                        else BlockType.CAPTION if is_caption
                        else BlockType.TEXT
                    )

                    result.blocks.append(
                        Block(
                            block_id=str(uuid.uuid4()),
                            doc_id=doc_id,
                            source_file=result.source_file,
                            block_type=block_type,
                            content=text,
                            page_number=slide_number,
                            order_index=order_index,
                        )
                    )
                    order_index += 1

            # Speaker notes, kept as text with a flag -- useful context but
            # not part of the visible slide flow.
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes_text:
                    result.blocks.append(
                        Block(
                            block_id=str(uuid.uuid4()),
                            doc_id=doc_id,
                            source_file=result.source_file,
                            block_type=BlockType.TEXT,
                            content=notes_text,
                            page_number=slide_number,
                            order_index=order_index,
                            extra_metadata={"is_speaker_note": True},
                        )
                    )
                    order_index += 1

        self._link_captions_to_media(result)
        return result

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _run_ocr(image_path: str, warnings: list) -> str:
        try:
            with Image.open(image_path) as im:
                return pytesseract.image_to_string(im)
        except Exception as e:
            warnings.append(f"OCR failed for {image_path}: {e}")
            return ""

    @staticmethod
    def _table_to_markdown(table):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not rows:
            return "", rows
        header, *body = rows
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * len(header)) + " |",
        ]
        for row in body:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines), rows

    @staticmethod
    def _link_captions_to_media(result: ExtractionResult):
        captions = [b for b in result.blocks if b.block_type == BlockType.CAPTION]
        media = [b for b in result.blocks if b.block_type in (BlockType.IMAGE, BlockType.TABLE)]

        for caption in captions:
            same_slide = [m for m in media if m.page_number == caption.page_number]
            nearby = [m for m in same_slide if abs(m.order_index - caption.order_index) <= CAPTION_LINK_WINDOW]
            if not nearby:
                continue
            nearest = min(nearby, key=lambda m: abs(m.order_index - caption.order_index))
            result.link(caption, nearest)